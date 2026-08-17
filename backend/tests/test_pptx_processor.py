from pathlib import Path

from pptx import Presentation

from app.pptx.processor import extract_presentation_content, translate_and_write_presentation
from app.services.translation import MockTranslationService


def test_extracts_text_boxes_and_placeholders(sample_pptx: Path) -> None:
    content = extract_presentation_content(sample_pptx)
    texts = [target.text for target in content.paragraph_targets]
    assert "Standalone text" in texts
    assert "AutoShape text" in texts
    assert "Quarterly Training" in texts
    assert "Policy review" in texts


def test_does_not_extract_table_url_footer(sample_pptx: Path) -> None:
    content = extract_presentation_content(sample_pptx)
    texts = [target.text for target in content.paragraph_targets]
    assert "Do not translate table" not in texts
    assert "www.example.com" not in texts


def test_reads_speaker_notes(sample_pptx: Path) -> None:
    content = extract_presentation_content(sample_pptx)
    assert content.notes[0].text == "Read this speaker note."


def test_no_notes_means_no_audio_targets(tmp_path: Path) -> None:
    path = tmp_path / "no-notes.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_textbox(0, 0, 1000, 1000).text = "Hello"
    prs.save(path)
    content = extract_presentation_content(path)
    assert content.notes == []


def test_translation_preserves_paragraph_order(sample_pptx: Path, tmp_path: Path) -> None:
    output = tmp_path / "translated.pptx"
    translate_and_write_presentation(
        sample_pptx,
        output,
        "en-US",
        "zh-TW",
        MockTranslationService(),
        lambda *_: None,
    )
    prs = Presentation(str(output))
    assert prs.slides[0].shapes.title.text.startswith("[MOCK_TRANSLATION:zh-TW] Quarterly Training")
    assert prs.slides[0].notes_slide.notes_text_frame.text.startswith("[MOCK_NARRATION_TRANSLATION:zh-TW]")


def test_output_pptx_can_be_opened(sample_pptx: Path, tmp_path: Path) -> None:
    output = tmp_path / "translated.pptx"
    translate_and_write_presentation(sample_pptx, output, "en-US", "th-TH", MockTranslationService(), lambda *_: None)
    opened = Presentation(str(output))
    assert len(opened.slides) == 3
