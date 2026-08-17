from pathlib import Path

import pytest
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches


@pytest.fixture()
def sample_pptx(tmp_path: Path) -> Path:
    path = tmp_path / "sample.pptx"
    prs = Presentation()
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = "Quarterly Training"
    title_slide.placeholders[1].text = "Security update"
    title_slide.notes_slide.notes_text_frame.text = "Read this speaker note."

    content_slide = prs.slides.add_slide(prs.slide_layouts[1])
    content_slide.shapes.title.text = "Agenda"
    content_slide.placeholders[1].text = "Policy review"

    blank_slide = prs.slides.add_slide(prs.slide_layouts[6])
    text_box = blank_slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1))
    text_box.text = "Standalone text"
    auto_shape = blank_slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.2), Inches(1), Inches(2), Inches(1))
    auto_shape.text = "AutoShape text"
    table = blank_slide.shapes.add_table(2, 2, Inches(1), Inches(2), Inches(4), Inches(1)).table
    table.cell(0, 0).text = "Do not translate table"
    footer = blank_slide.shapes.add_textbox(Inches(1), Inches(6), Inches(3), Inches(0.4))
    footer.name = "Footer"
    footer.text = "www.example.com"

    prs.save(path)
    return path
